#!/usr/bin/env python3
"""
MailMind — AlgoQuest 2025 Round 2 Presentation Generator
Team Cipher | AI-Powered Smart Email Assistant

Generates: mailmindd/MailMind_AlgoQuest_R2.pptx
Run:       python mailmindd/generate_ppt.py
Requires:  pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ──────────────────────────────────────────────────────────────
BG_COLOR        = RGBColor(15, 23, 42)       # #0F172A  deep navy
WHITE           = RGBColor(255, 255, 255)
LIGHT_GRAY      = RGBColor(148, 163, 184)    # #94A3B8
MID_GRAY        = RGBColor(100, 116, 139)    # #64748B
ELECTRIC_BLUE   = RGBColor(59, 130, 246)     # #3B82F6
TEAL            = RGBColor(6, 182, 212)      # #06B6D4
PURPLE          = RGBColor(139, 92, 246)     # #8B5CF6
RED             = RGBColor(239, 68, 68)      # #EF4444
GREEN           = RGBColor(34, 197, 94)      # #22C55E
AMBER           = RGBColor(245, 158, 11)     # #F59E0B
CARD_BG         = RGBColor(30, 41, 59)       # #1E293B  slightly lighter navy
CARD_BG_LIGHT   = RGBColor(51, 65, 85)       # #334155

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT    = "Calibri"

# ──────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ──────────────────────────────────────────────────────────────

def _no_border(shape):
    """Remove the outline / border from a shape."""
    shape.line.fill.background()


def add_background(slide):
    """Fill the entire slide with the dark navy background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_accent_bar(slide, color=ELECTRIC_BLUE, height=Inches(0.08)):
    """Add a thin coloured bar across the very top of the slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    _no_border(bar)
    return bar


def add_footer(slide, text="Team Cipher  |  AlgoQuest 2025"):
    """Add a small grey footer at the bottom-right."""
    tb = slide.shapes.add_textbox(
        Inches(8.5), Inches(7.05), Inches(4.5), Inches(0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MID_GRAY
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT
    _no_border(tb)


def add_text_box(slide, left, top, width, height, text,
                 font_size=Pt(18), color=WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name=FONT,
                 word_wrap=True):
    """Convenience: add a simple single-paragraph text box."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    _no_border(tb)
    return tb


def add_rich_text_box(slide, left, top, width, height):
    """Return (text_frame, textbox) so caller can add multiple paragraphs."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    _no_border(tb)
    return tf, tb


def _add_run(paragraph, text, size=Pt(18), color=WHITE, bold=False, name=FONT):
    """Add a run to an existing paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    return run


def _add_paragraph(tf, text="", size=Pt(18), color=WHITE, bold=False,
                   alignment=PP_ALIGN.LEFT, space_after=Pt(6), name=FONT):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = name
    p.alignment = alignment
    p.space_after = space_after
    return p


def add_card(slide, left, top, width, height, fill_color=CARD_BG):
    """Add a filled rounded-look rectangle (card) and return the shape."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    _no_border(card)
    return card


def add_rect(slide, left, top, width, height, fill_color=ELECTRIC_BLUE):
    """Add a plain rectangle shape."""
    r = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    r.fill.solid()
    r.fill.fore_color.rgb = fill_color
    _no_border(r)
    return r


def add_slide_title(slide, title, subtitle=None):
    """Add a large title (and optional subtitle) near the top of the slide."""
    add_text_box(slide, 0.8, 0.35, 11.5, 0.7, title,
                 font_size=Pt(38), color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, 0.8, 1.0, 11.5, 0.5, subtitle,
                     font_size=Pt(18), color=LIGHT_GRAY, bold=False,
                     alignment=PP_ALIGN.LEFT)


def new_slide(prs):
    """Create a blank slide with background, accent bar and footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_accent_bar(slide)
    add_footer(slide)
    return slide


# ──────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ──────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide."""
    slide = new_slide(prs)

    # Decorative accent line in centre
    add_rect(slide, 4.5, 1.2, 4.3, 0.06, ELECTRIC_BLUE)

    # Main title
    add_text_box(slide, 1, 1.5, 11.3, 1.0,
                 "\U0001f9e0  MailMind",                 # 🧠
                 font_size=Pt(54), color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, 1, 2.6, 11.3, 0.6,
                 "AI-Powered Smart Email Assistant",
                 font_size=Pt(28), color=ELECTRIC_BLUE, bold=False,
                 alignment=PP_ALIGN.CENTER)

    # Problem statement
    add_text_box(slide, 1, 3.35, 11.3, 0.5,
                 "Problem Statement #1 — Smart Email Solutions",
                 font_size=Pt(18), color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Team
    add_text_box(slide, 1, 4.1, 11.3, 0.5,
                 "Team Cipher  |  AlgoQuest 2025 — Round 2",
                 font_size=Pt(20), color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Decorative line
    add_rect(slide, 3.5, 4.85, 6.3, 0.04, CARD_BG_LIGHT)

    # Tech badges row
    badges = "Next.js 16  •  React 19  •  TypeScript 5  •  Groq Llama 3.3  •  Vitest"
    add_text_box(slide, 1, 5.1, 11.3, 0.5, badges,
                 font_size=Pt(16), color=TEAL,
                 alignment=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """The Problem: Email Overload."""
    slide = new_slide(prs)
    add_slide_title(slide, "The Problem: Email Overload",
                    "Professionals are drowning — and current tools aren't helping.")

    stats = [
        ("120+",  "emails / day",       "Professionals are drowning in their inbox",   ELECTRIC_BLUE),
        ("28%",   "of work time",       "Spent just managing email",                   AMBER),
        ("Missed","deadlines",          "Critical dates buried in email text",          RED),
        ("No",    "context awareness",  "Existing tools just match keywords",           PURPLE),
    ]

    for i, (big, label, desc, color) in enumerate(stats):
        x = 0.6 + i * 3.1
        # Card background
        add_card(slide, x, 1.8, 2.8, 4.5)

        # Big stat
        add_text_box(slide, x + 0.2, 2.1, 2.4, 0.9, big,
                     font_size=Pt(48), color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        # Label
        add_text_box(slide, x + 0.2, 3.0, 2.4, 0.5, label,
                     font_size=Pt(20), color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        # Divider
        add_rect(slide, x + 0.6, 3.55, 1.6, 0.04, color)
        # Description
        add_text_box(slide, x + 0.25, 3.8, 2.3, 1.5, desc,
                     font_size=Pt(15), color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)


def slide_03_solution(prs):
    """Our Solution: MailMind."""
    slide = new_slide(prs)
    add_slide_title(slide, "Our Solution: MailMind")

    # One-liner
    add_text_box(slide, 0.8, 1.2, 11.5, 0.8,
                 "An AI-native email assistant powered by Groq's Llama 3.3 70B\n"
                 "that understands your emails contextually — not just keyword matching.",
                 font_size=Pt(20), color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    cards_data = [
        ("\U0001f9e0  Contextual AI",
         "Understands nuance, intent, and urgency — not just keywords. "
         "Powered by Groq Llama 3.3 70B Versatile.",
         ELECTRIC_BLUE),
        ("\U0001f916  Agentic Automation",
         "One-click \"Handle For Me\" — the AI reads, categorises, drafts a reply, "
         "extracts events, and creates tasks autonomously.",
         TEAL),
        ("\U0001f9ea  LLM-Validated Testing",
         "AI tests AI: an independent LLM oracle validates every AI output "
         "for correctness, quality, and confidence.",
         PURPLE),
    ]

    for i, (title, body, color) in enumerate(cards_data):
        x = 0.6 + i * 4.1
        add_card(slide, x, 2.5, 3.8, 4.0)
        # Color accent bar on top of card
        add_rect(slide, x, 2.5, 3.8, 0.08, color)

        add_text_box(slide, x + 0.25, 2.8, 3.3, 0.6, title,
                     font_size=Pt(22), color=color, bold=True,
                     alignment=PP_ALIGN.LEFT)

        add_text_box(slide, x + 0.25, 3.5, 3.3, 2.8, body,
                     font_size=Pt(16), color=LIGHT_GRAY,
                     alignment=PP_ALIGN.LEFT)


def slide_04_mapping(prs):
    """Problem Statement → Our Implementation."""
    slide = new_slide(prs)
    add_slide_title(slide, "Problem Statement → Our Implementation",
                    "Every requirement mapped to a concrete feature.")

    rows = [
        ("Auto-prioritize mails",    "AI Priority Scoring (1-100)",              ELECTRIC_BLUE),
        ("Extract tasks",            "AI To-Do Extraction + Agentic Handle For Me", TEAL),
        ("Manage follow-ups",        "AI Follow-Up Scheduling",                  GREEN),
        ("NLP",                      "Deadline Extraction, Summarization, Spam Detection", PURPLE),
        ("RAG",                      "In-Memory Vector Store + Cosine Similarity Replies", AMBER),
        ("React",                    "React 19 + Next.js 16 + Custom Hooks",     ELECTRIC_BLUE),
        ("Node.js",                  "16+ Serverless API Routes",                TEAL),
        ("Testing with LLMs",        "LLM-as-Test-Oracle (5 Test Suites)",       RED),
    ]

    y_start = 1.7
    row_h = 0.62

    # Header
    add_rect(slide, 0.6, y_start, 5.4, row_h, CARD_BG_LIGHT)
    add_text_box(slide, 0.8, y_start + 0.1, 5.0, 0.4,
                 "Requirement", font_size=Pt(16), color=ELECTRIC_BLUE, bold=True)

    add_rect(slide, 6.1, y_start, 6.6, row_h, CARD_BG_LIGHT)
    add_text_box(slide, 6.3, y_start + 0.1, 6.2, 0.4,
                 "MailMind Implementation", font_size=Pt(16), color=ELECTRIC_BLUE, bold=True)

    for idx, (req, impl, color) in enumerate(rows):
        y = y_start + (idx + 1) * row_h
        bg = CARD_BG if idx % 2 == 0 else RGBColor(22, 33, 50)

        # Left cell
        add_rect(slide, 0.6, y, 5.4, row_h, bg)
        # Color dot
        add_rect(slide, 0.6, y, 0.08, row_h, color)
        add_text_box(slide, 0.85, y + 0.1, 5.0, 0.4, req,
                     font_size=Pt(15), color=WHITE, bold=True)

        # Right cell
        add_rect(slide, 6.1, y, 6.6, row_h, bg)
        add_text_box(slide, 6.3, y + 0.1, 6.2, 0.4, impl,
                     font_size=Pt(15), color=LIGHT_GRAY)


def slide_05_architecture(prs):
    """System Architecture."""
    slide = new_slide(prs)
    add_slide_title(slide, "System Architecture",
                    "Clean, serverless, AI-first design.")

    # ── Frontend Layer ──
    add_card(slide, 0.8, 1.6, 11.7, 1.0, CARD_BG)
    add_rect(slide, 0.8, 1.6, 11.7, 0.07, ELECTRIC_BLUE)
    add_text_box(slide, 1.0, 1.75, 11.3, 0.35,
                 "\U0001f5a5  Frontend Layer",
                 font_size=Pt(18), color=ELECTRIC_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.0, 2.1, 11.3, 0.35,
                 "Next.js 16  •  React 19  •  TypeScript 5  •  Tailwind CSS  •  Custom Hooks",
                 font_size=Pt(14), color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, 6.0, 2.65, 1.3, 0.4, "▼",
                 font_size=Pt(24), color=ELECTRIC_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # ── API Layer ──
    add_card(slide, 0.8, 3.0, 11.7, 0.7, CARD_BG)
    add_rect(slide, 0.8, 3.0, 11.7, 0.07, TEAL)
    add_text_box(slide, 1.0, 3.1, 11.3, 0.5,
                 "\U0001f517  API Layer  —  16+ Serverless Endpoints",
                 font_size=Pt(18), color=TEAL, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, 6.0, 3.7, 1.3, 0.4, "▼",
                 font_size=Pt(24), color=TEAL,
                 alignment=PP_ALIGN.CENTER)

    # ── Sub-modules ──
    modules = [
        ("AI Engine\n10 endpoints", PURPLE),
        ("Gmail API\nIntegration",  ELECTRIC_BLUE),
        ("RAG\nVector Store",       TEAL),
        ("Calendar\nExtraction",    AMBER),
        ("Team\nCollab",            GREEN),
        ("Search &\nFilter",        RED),
    ]
    for i, (label, color) in enumerate(modules):
        x = 0.8 + i * 2.05
        add_card(slide, x, 4.1, 1.85, 1.3, CARD_BG)
        add_rect(slide, x, 4.1, 1.85, 0.06, color)
        add_text_box(slide, x + 0.1, 4.25, 1.65, 1.0, label,
                     font_size=Pt(13), color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Arrow down
    add_text_box(slide, 6.0, 5.4, 1.3, 0.4, "▼",
                 font_size=Pt(24), color=PURPLE,
                 alignment=PP_ALIGN.CENTER)

    # ── External Services ──
    add_card(slide, 0.8, 5.7, 11.7, 1.0, CARD_BG)
    add_rect(slide, 0.8, 5.7, 11.7, 0.07, PURPLE)
    add_text_box(slide, 1.0, 5.85, 11.3, 0.35,
                 "\u2601  External Services",
                 font_size=Pt(18), color=PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.0, 6.2, 11.3, 0.35,
                 "Groq Llama 3.3 70B Versatile  •  Gmail API  •  NextAuth OAuth 2.0",
                 font_size=Pt(14), color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)


def slide_06_core_features(prs):
    """Core AI Features."""
    slide = new_slide(prs)
    add_slide_title(slide, "Core AI Features",
                    "Intelligent email processing powered by Groq Llama 3.3 70B.")

    features = [
        ("\U0001f3af  Priority Scoring (1-100)",
         "AI analyses urgency, sender importance, and action requirements. "
         "Returns a numeric score with reasoning — no more guessing "
         "which emails matter most.",
         ELECTRIC_BLUE),
        ("\U0001f4c2  Smart 4-Category Inbox",
         "Do Now  |  Needs Decision  |  Waiting  |  Low Energy\n\n"
         "Emails are auto-sorted by the AI into actionable buckets "
         "so you always know what to tackle first.",
         TEAL),
        ("\U0001f6e1  Spam & Phishing Detection",
         "Context-aware detection with confidence scoring. "
         "Catches subtle attacks like paypa1.com vs paypal.com — "
         "beyond what rule-based filters can do.",
         RED),
    ]

    for i, (title, body, color) in enumerate(features):
        x = 0.6 + i * 4.1
        add_card(slide, x, 1.8, 3.8, 4.8)
        add_rect(slide, x, 1.8, 3.8, 0.08, color)

        add_text_box(slide, x + 0.25, 2.1, 3.3, 0.6, title,
                     font_size=Pt(20), color=color, bold=True)

        add_text_box(slide, x + 0.25, 2.8, 3.3, 3.5, body,
                     font_size=Pt(15), color=LIGHT_GRAY)


def slide_07_nlp_rag(prs):
    """NLP & RAG-Powered Intelligence."""
    slide = new_slide(prs)
    add_slide_title(slide, "NLP & RAG-Powered Intelligence",
                    "Deep language understanding meets retrieval-augmented generation.")

    # ── Left: NLP ──
    add_card(slide, 0.6, 1.7, 5.8, 5.0)
    add_rect(slide, 0.6, 1.7, 5.8, 0.08, ELECTRIC_BLUE)
    add_text_box(slide, 0.9, 1.95, 5.2, 0.5,
                 "\U0001f4ac  NLP Features", font_size=Pt(22),
                 color=ELECTRIC_BLUE, bold=True)

    nlp_items = [
        ("Deadline Extraction", "\"by end of week\" → Fri 2025-01-31"),
        ("AI Email Summarization", "Long threads → concise bullet points"),
        ("Sentiment & Intent", "Detects urgency, frustration, requests"),
        ("\"Why This Matters\"", "AI explains why each email needs attention"),
    ]
    y = 2.6
    for title, desc in nlp_items:
        add_rect(slide, 1.0, y, 0.08, 0.55, ELECTRIC_BLUE)
        add_text_box(slide, 1.25, y, 4.9, 0.28, title,
                     font_size=Pt(16), color=WHITE, bold=True)
        add_text_box(slide, 1.25, y + 0.3, 4.9, 0.28, desc,
                     font_size=Pt(13), color=LIGHT_GRAY)
        y += 0.95

    # ── Right: RAG ──
    add_card(slide, 6.9, 1.7, 5.8, 5.0)
    add_rect(slide, 6.9, 1.7, 5.8, 0.08, TEAL)
    add_text_box(slide, 7.2, 1.95, 5.2, 0.5,
                 "\U0001f50d  RAG Implementation", font_size=Pt(22),
                 color=TEAL, bold=True)

    rag_items = [
        ("In-Memory Vector Store", "128-dim TF-IDF embeddings"),
        ("Cosine Similarity Matching", "Finds relevant past emails instantly"),
        ("Sender-Aware Retrieval", "Prioritises context from same sender"),
        ("500 Email Capacity", "Stores up to 500 email embeddings"),
        ("Context-Enriched Replies", "RAG-powered smart reply generation"),
    ]
    y = 2.6
    for title, desc in rag_items:
        add_rect(slide, 7.3, y, 0.08, 0.55, TEAL)
        add_text_box(slide, 7.55, y, 4.9, 0.28, title,
                     font_size=Pt(16), color=WHITE, bold=True)
        add_text_box(slide, 7.55, y + 0.3, 4.9, 0.28, desc,
                     font_size=Pt(13), color=LIGHT_GRAY)
        y += 0.85


def slide_08_agentic(prs):
    """Agentic AI: One-Click Email Handling."""
    slide = new_slide(prs)
    add_slide_title(slide, "Agentic AI: One-Click Email Handling",
                    "The AI handles everything — autonomously, step by step.")

    steps = [
        ("1", "Analyze",       "Summarises\nthe email",          ELECTRIC_BLUE),
        ("2", "Categorize",    "Determines\npriority & category", TEAL),
        ("3", "Draft Reply",   "Generates context-\naware response", PURPLE),
        ("4", "Extract Events","Identifies\ncalendar events",     AMBER),
        ("5", "Generate Task", "Creates\nactionable to-do",       GREEN),
        ("6", "Follow-Up",    "Recommends\nnext steps",           RED),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = 0.5 + i * 2.1

        # Card
        add_card(slide, x, 2.0, 1.9, 3.8)
        add_rect(slide, x, 2.0, 1.9, 0.08, color)

        # Step number circle
        add_text_box(slide, x + 0.55, 2.25, 0.8, 0.6, num,
                     font_size=Pt(32), color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + 0.1, 3.0, 1.7, 0.5, title,
                     font_size=Pt(17), color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

        # Desc
        add_text_box(slide, x + 0.1, 3.55, 1.7, 1.5, desc,
                     font_size=Pt(13), color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

        # Arrow between cards
        if i < 5:
            add_text_box(slide, x + 1.85, 3.0, 0.35, 0.5, "→",
                         font_size=Pt(24), color=MID_GRAY,
                         alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_card(slide, 0.8, 6.1, 11.7, 0.7, CARD_BG)
    add_text_box(slide, 1.0, 6.2, 11.3, 0.5,
                 "All autonomous.  One click.  Step-by-step progress modal.  "
                 "No manual intervention required.",
                 font_size=Pt(16), color=TEAL, alignment=PP_ALIGN.CENTER)


def slide_09_productivity(prs):
    """Productivity & Collaboration Suite."""
    slide = new_slide(prs)
    add_slide_title(slide, "Productivity & Collaboration Suite",
                    "Beyond email — a complete workflow platform.")

    cards = [
        # row 1
        ("\U0001f3af  Focus Mode",
         "Distraction-free urgent task view with time-of-day greeting, "
         "progress tracking, and smart task prioritisation.",
         ELECTRIC_BLUE, 0.6, 1.8),
        ("\U0001f4ca  Weekly Analysis",
         "Email volume tracking, stress scoring (0-100), burnout risk "
         "detection (Low → Critical), late-night email flags.",
         AMBER, 6.6, 1.8),
        # row 2
        ("\U0001f4c5  Calendar Integration",
         "AI extracts events from emails. Color-coded: "
         "Deadline (red), Meeting (blue), Appointment (purple), Reminder (green).",
         GREEN, 0.6, 4.2),
        ("\U0001f465  Team Collaboration",
         "Email assignment, workload dashboard, status tracking, "
         "internal notes, AI workload suggestions.",
         PURPLE, 6.6, 4.2),
    ]

    for title, body, color, x, y in cards:
        add_card(slide, x, y, 5.7, 2.1)
        add_rect(slide, x, y, 5.7, 0.08, color)
        add_text_box(slide, x + 0.25, y + 0.2, 5.2, 0.5, title,
                     font_size=Pt(22), color=color, bold=True)
        add_text_box(slide, x + 0.25, y + 0.8, 5.2, 1.1, body,
                     font_size=Pt(15), color=LIGHT_GRAY)


def slide_10_testing(prs):
    """Innovation: LLM-as-Test-Oracle."""
    slide = new_slide(prs)
    add_slide_title(slide, "Innovation: LLM-as-Test-Oracle",
                    "How do you test if AI output is reasonable?  You ask another LLM.")

    # Visual flow
    flow_items = [
        ("Test\nInput",        CARD_BG_LIGHT),
        ("→",                  None),
        ("MailMind AI\n(Groq)", ELECTRIC_BLUE),
        ("→",                  None),
        ("AI Output",          CARD_BG_LIGHT),
        ("→",                  None),
        ("LLM Oracle\n(Groq)", PURPLE),
        ("→",                  None),
        ("Validation\n✓ / ✗",  GREEN),
    ]
    x = 0.3
    for label, color in flow_items:
        if color is None:
            add_text_box(slide, x, 1.95, 0.5, 0.7, "→",
                         font_size=Pt(28), color=MID_GRAY,
                         alignment=PP_ALIGN.CENTER)
            x += 0.45
        else:
            add_card(slide, x, 1.8, 1.35, 1.0, color)
            add_text_box(slide, x + 0.05, 1.85, 1.25, 0.85, label,
                         font_size=Pt(12), color=WHITE, bold=True,
                         alignment=PP_ALIGN.CENTER)
            x += 1.45

    # Validation output box
    add_card(slide, 8.5, 3.0, 4.2, 0.8, CARD_BG)
    add_text_box(slide, 8.7, 3.1, 3.8, 0.6,
                 "{ isValid, confidence, reasoning }",
                 font_size=Pt(14), color=GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Test suites table
    suites = [
        ("ai-priority.test.ts",           "Validates priority scores (1-100)",        ELECTRIC_BLUE),
        ("ai-categorization.test.ts",      "Validates category assignments",           TEAL),
        ("ai-spam-detection.test.ts",      "Validates spam / phishing decisions",      RED),
        ("ai-deadline-extraction.test.ts", "Validates deadline parsing accuracy",      AMBER),
        ("ai-reply-quality.test.ts",       "Validates reply professionalism & tone",   PURPLE),
    ]

    y = 4.0
    for name, desc, color in suites:
        add_rect(slide, 0.8, y, 0.08, 0.55, color)
        add_card(slide, 0.95, y, 11.5, 0.55, CARD_BG)
        add_text_box(slide, 1.1, y + 0.08, 4.5, 0.4, name,
                     font_size=Pt(15), color=color, bold=True)
        add_text_box(slide, 5.8, y + 0.08, 6.0, 0.4, desc,
                     font_size=Pt(15), color=LIGHT_GRAY)
        y += 0.65

    # Bottom note
    add_text_box(slide, 0.8, 7.0 - 0.6, 11.5, 0.4,
                 "Powered by Vitest  +  Groq Llama 3.3 70B Versatile",
                 font_size=Pt(14), color=TEAL, alignment=PP_ALIGN.CENTER)


def slide_11_scalability(prs):
    """Real-World Ready."""
    slide = new_slide(prs)
    add_slide_title(slide, "Real-World Ready",
                    "Built for production from day one.")

    metrics = [
        ("<2s",   "AI Response Time",        "Groq's ultra-fast inference", ELECTRIC_BLUE),
        ("16+",   "Serverless Endpoints",    "Modular API design",          TEAL),
        ("500",   "Emails in RAG Store",     "In-memory vector DB",         PURPLE),
        ("10",    "Parallel Emails",         "Concurrent processing",       AMBER),
    ]

    for i, (big, label, sub, color) in enumerate(metrics):
        x = 0.5 + i * 3.2
        add_card(slide, x, 1.7, 2.9, 2.8)
        add_rect(slide, x, 1.7, 2.9, 0.08, color)

        add_text_box(slide, x + 0.15, 1.95, 2.6, 0.8, big,
                     font_size=Pt(52), color=color, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, 2.85, 2.6, 0.4, label,
                     font_size=Pt(18), color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, 3.3, 2.6, 0.4, sub,
                     font_size=Pt(13), color=LIGHT_GRAY,
                     alignment=PP_ALIGN.CENTER)

    # Scalability points
    points = [
        ("Serverless Architecture",  "Auto-scales with demand — zero ops overhead",      ELECTRIC_BLUE),
        ("In-Memory Caching",        "Reduces redundant API calls for instant responses", TEAL),
        ("Modular API Design",       "Each feature is an independent, testable endpoint", PURPLE),
        ("Real Gmail OAuth",         "Production-ready NextAuth integration",             GREEN),
    ]

    y = 4.9
    for title, desc, color in points:
        add_rect(slide, 0.8, y, 0.08, 0.5, color)
        add_text_box(slide, 1.1, y + 0.05, 4.0, 0.4, title,
                     font_size=Pt(17), color=WHITE, bold=True)
        add_text_box(slide, 5.3, y + 0.05, 7.2, 0.4, desc,
                     font_size=Pt(15), color=LIGHT_GRAY)
        y += 0.6


def slide_12_thanks(prs):
    """Thank You + Live Demo."""
    slide = new_slide(prs)

    # Decorative centre line
    add_rect(slide, 4.5, 1.0, 4.3, 0.06, ELECTRIC_BLUE)

    # Big thank you
    add_text_box(slide, 1, 1.3, 11.3, 1.0,
                 "Thank You!",
                 font_size=Pt(52), color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, 1, 2.4, 11.3, 0.6,
                 "Let us show you MailMind in action…",
                 font_size=Pt(24), color=ELECTRIC_BLUE,
                 alignment=PP_ALIGN.CENTER)

    add_rect(slide, 4.5, 3.2, 4.3, 0.04, CARD_BG_LIGHT)

    # Team
    add_text_box(slide, 1, 3.5, 11.3, 0.5,
                 "Team Cipher",
                 font_size=Pt(28), color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # GitHub
    add_text_box(slide, 1, 4.2, 11.3, 0.4,
                 "GitHub:  github.com/Avila-Princy-M01/mailmindd",
                 font_size=Pt(16), color=TEAL,
                 alignment=PP_ALIGN.CENTER)

    # Tech stack
    add_text_box(slide, 1, 4.8, 11.3, 0.4,
                 "Next.js 16  •  React 19  •  Groq Llama 3.3 70B  •  Vitest",
                 font_size=Pt(16), color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)

    # Questions
    add_card(slide, 4.5, 5.5, 4.3, 1.0, CARD_BG)
    add_rect(slide, 4.5, 5.5, 4.3, 0.08, PURPLE)
    add_text_box(slide, 4.7, 5.7, 3.9, 0.6,
                 "Questions?",
                 font_size=Pt(30), color=PURPLE, bold=True,
                 alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────

def generate():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_mapping(prs)
    slide_05_architecture(prs)
    slide_06_core_features(prs)
    slide_07_nlp_rag(prs)
    slide_08_agentic(prs)
    slide_09_productivity(prs)
    slide_10_testing(prs)
    slide_11_scalability(prs)
    slide_12_thanks(prs)

    # Determine output path relative to this script's directory
    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(script_dir, "MailMind_AlgoQuest_R2.pptx")
    prs.save(out_path)
    print("[OK] Presentation saved -> {}".format(out_path))
    print("     Slides: {}".format(len(prs.slides)))


if __name__ == "__main__":
    generate()
